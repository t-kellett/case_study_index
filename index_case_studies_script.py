import os
from pptx import Presentation
from openpyxl import Workbook, load_workbook

# Define the paths to the PowerPoint presentation and the Excel file
pptx_path = "/Users/tom.kellett/Internal/case_study_index/test/test_case_studies_1.pptx"
excel_path = "/Users/tom.kellett/Internal/case_study_index/Case Study Index Test.xlsx"

# Create or open the Excel workbook
if os.path.exists(excel_path):
    wb = load_workbook(excel_path)
else:
    wb = Workbook()

# Create or select the 'Case Studies' worksheet
if 'Case Studies' in wb.sheetnames:
    ws = wb['Case Studies']
else:
    ws = wb.create_sheet('Case Studies')

# Add or update column headings
column_headings = ['Client', 'Industry', 'Service Line', 'Offering', 'Sub Offering', 'Team', 'Notes', 'Slide Number']
for col_num, heading in enumerate(column_headings, start=1):
    cell = ws.cell(row=1, column=col_num, value=heading)

# Load the PowerPoint presentation
presentation = Presentation(pptx_path)

# Iterate through the slides and extract data from speaker notes
for slide_number, slide in enumerate(presentation.slides, start=1):
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame.text

    if notes_text:
        data = dict()
        lines = notes_text.split('\n')
        
        # Process each line and handle multiline values
        for line in lines:
            parts = line.split(': ', 1)
            if len(parts) == 2:
                key, value = map(str.strip, parts)
                if key in data:
                    data[key] += ', ' + value
                else:
                    data[key] = value

        # Set the "Slide Number" value in the data dictionary
        data['Slide Number'] = slide_number

        # Populate Excel rows with data from speaker notes
        for col_num, key in enumerate(column_headings, start=1):
            cell = ws.cell(row=slide_number + 1, column=col_num, value=data.get(key, ''))

# Save and close the Excel workbook
wb.save(excel_path)
wb.close()
